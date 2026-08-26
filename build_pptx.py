#!/usr/bin/env python3
"""Сборка редактируемого PPTX по актуальному HTML-сценарию + анимации появления по клику."""

from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

ROOT = Path(__file__).resolve().parent
OUT = ROOT / "Лингвистический взлом — Лучший лектор РУДН.pptx"

W, H = Inches(13.333), Inches(7.5)
PAPER = RGBColor(0xF3, 0xEB, 0xE0)
WHITE = RGBColor(0xFF, 0xFA, 0xF3)
INK = RGBColor(0x2A, 0x22, 0x18)
MUTED = RGBColor(0x6B, 0x5E, 0x4F)
BURG = RGBColor(0x7A, 0x1F, 0x2B)
NAVY = RGBColor(0x1A, 0x27, 0x44)
GOLD = RGBColor(0xC4, 0xA3, 0x5A)
FR = RGBColor(0x3D, 0x5A, 0x80)
RU = RGBColor(0x8E, 0x24, 0x28)
CREAM = RGBColor(0xE8, 0xD5, 0xA3)
DARK = RGBColor(0x0D, 0x12, 0x1C)
P = "http://schemas.openxmlformats.org/presentationml/2006/main"


def set_run(run, text, size=18, bold=False, italic=False, color=INK, font="Calibri"):
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font


def add_text(slide, l, t, w, h, text, size=18, bold=False, italic=False, color=INK,
             font="Calibri", align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    try:
        tf._txBody.bodyPr.set("anchor", {MSO_ANCHOR.TOP: "t", MSO_ANCHOR.MIDDLE: "ctr", MSO_ANCHOR.BOTTOM: "b"}[anchor])
    except Exception:
        pass
    for i, line in enumerate(str(text).split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        set_run(run, line, size, bold, italic, color, font)
    return box


def p_run(paragraph, text, size=18, bold=False, italic=False, color=INK, font="Calibri"):
    run = paragraph.add_run()
    set_run(run, text, size, bold, italic, color, font)
    return run


def fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def rect(slide, l, t, w, h, color):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    fill(s, color)
    return s


def card(slide, l, t, w, h, color=WHITE):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    fill(s, color)
    s.line.color.rgb = RGBColor(0xE0, 0xD2, 0xB8)
    try:
        s.adjustments[0] = 0.04
    except Exception:
        pass
    return s


def eyebrow(slide, text, l=Inches(0.7), t=Inches(0.32), color=GOLD):
    return add_text(slide, l, t, Inches(12), Inches(0.35), text.upper(), 12, True, False, color)


def heading(slide, text, t=Inches(0.58), size=36, color=INK):
    return add_text(slide, Inches(0.7), t, Inches(12), Inches(0.75), text, size, True, False, color, "Georgia")


def paper_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    rect(slide, 0, 0, W, H, PAPER)
    rect(slide, 0, H - Inches(0.06), W, Inches(0.06), GOLD)
    return slide


def dark_slide(prs, color=DARK):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    rect(slide, 0, 0, W, H, color)
    return slide


def chip(slide, text, l, t, color=WHITE, ink=BURG):
    width = Inches(max(1.2, 0.26 * len(text) + 0.5))
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, width, Inches(0.5))
    fill(s, color)
    s.line.color.rgb = GOLD
    tf = s.text_frame
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    p_run(tf.paragraphs[0], text, 16, False, False, ink, "Georgia")
    return s, width


def bag_items(slide, words, l, t, extra=False, max_r=None, shapes_out=None):
    x, y = l, t
    if max_r is None:
        max_r = Inches(12.6)
    for w in words:
        width = Inches(max(1.1, 0.22 * len(w) + 0.38))
        if x + width > max_r:
            x = l
            y += Inches(0.48)
        s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, width, Inches(0.4))
        fill(s, RGBColor(0xF7, 0xE6, 0xE4) if extra else WHITE)
        s.line.color.rgb = RGBColor(0xC4, 0x8A, 0x86) if extra else RGBColor(0xE0, 0xD2, 0xB8)
        tf = s.text_frame
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        p_run(tf.paragraphs[0], w, 13, False, False, RU if extra else INK)
        if shapes_out is not None:
            shapes_out.append(s)
        x += width + Inches(0.08)
    return y


def _shape_id(shape):
    return str(shape.shape_id)


def add_click_appears(slide, shape_groups):
    """Каждая группа появляется по клику (fade + visibility)."""
    if not shape_groups:
        return
    normalized = []
    for g in shape_groups:
        if isinstance(g, (list, tuple)):
            normalized.append([x for x in g if x is not None])
        elif g is not None:
            normalized.append([g])
    normalized = [g for g in normalized if g]
    if not normalized:
        return

    cSld = slide.shapes._spTree.getparent()
    # remove old timing if any
    for old in cSld.findall(qn("p:timing")):
        cSld.remove(old)

    timing = etree.SubElement(cSld, qn("p:timing"))
    tnLst = etree.SubElement(timing, qn("p:tnLst"))
    root_par = etree.SubElement(tnLst, qn("p:par"))
    root_cTn = etree.SubElement(root_par, qn("p:cTn"))
    root_cTn.set("id", "1")
    root_cTn.set("dur", "indefinite")
    root_cTn.set("restart", "never")
    root_cTn.set("nodeType", "tmRoot")
    root_child = etree.SubElement(root_cTn, qn("p:childTnLst"))

    seq = etree.SubElement(root_child, qn("p:seq"))
    seq.set("concurrent", "true")
    seq.set("nextAc", "seek")
    main_cTn = etree.SubElement(seq, qn("p:cTn"))
    main_cTn.set("id", "2")
    main_cTn.set("dur", "indefinite")
    main_cTn.set("nodeType", "mainSeq")
    main_child = etree.SubElement(main_cTn, qn("p:childTnLst"))

    next_id = 3
    for group in normalized:
        click_par = etree.SubElement(main_child, qn("p:par"))
        click_cTn = etree.SubElement(click_par, qn("p:cTn"))
        click_cTn.set("id", str(next_id)); next_id += 1
        click_cTn.set("fill", "hold")
        click_cTn.set("nodeType", "clickEffect")
        st = etree.SubElement(click_cTn, qn("p:stCondLst"))
        cond = etree.SubElement(st, qn("p:cond"))
        cond.set("delay", "0")
        click_child = etree.SubElement(click_cTn, qn("p:childTnLst"))

        for shape in group:
            sid = _shape_id(shape)
            inner_par = etree.SubElement(click_child, qn("p:par"))
            inner_cTn = etree.SubElement(inner_par, qn("p:cTn"))
            inner_cTn.set("id", str(next_id)); next_id += 1
            inner_cTn.set("fill", "hold")
            st2 = etree.SubElement(inner_cTn, qn("p:stCondLst"))
            cond2 = etree.SubElement(st2, qn("p:cond"))
            cond2.set("delay", "0")
            effects = etree.SubElement(inner_cTn, qn("p:childTnLst"))

            # set visibility
            set_el = etree.SubElement(effects, qn("p:set"))
            cBhvr = etree.SubElement(set_el, qn("p:cBhvr"))
            cTn = etree.SubElement(cBhvr, qn("p:cTn"))
            cTn.set("id", str(next_id)); next_id += 1
            cTn.set("dur", "1")
            cTn.set("fill", "hold")
            st3 = etree.SubElement(cTn, qn("p:stCondLst"))
            cond3 = etree.SubElement(st3, qn("p:cond"))
            cond3.set("delay", "0")
            tgt = etree.SubElement(cBhvr, qn("p:tgtEl"))
            spTgt = etree.SubElement(tgt, qn("p:spTgt"))
            spTgt.set("spid", sid)
            attr = etree.SubElement(cBhvr, qn("p:attrNameLst"))
            an = etree.SubElement(attr, qn("p:attrName"))
            an.text = "style.visibility"
            to = etree.SubElement(set_el, qn("p:to"))
            sval = etree.SubElement(to, qn("p:strVal"))
            sval.set("val", "visible")

            # fade in
            anim = etree.SubElement(effects, qn("p:animEffect"))
            anim.set("transition", "in")
            anim.set("filter", "fade")
            cBhvr2 = etree.SubElement(anim, qn("p:cBhvr"))
            cTn2 = etree.SubElement(cBhvr2, qn("p:cTn"))
            cTn2.set("id", str(next_id)); next_id += 1
            cTn2.set("dur", "450")
            tgt2 = etree.SubElement(cBhvr2, qn("p:tgtEl"))
            spTgt2 = etree.SubElement(tgt2, qn("p:spTgt"))
            spTgt2.set("spid", sid)

    # prevCondLst / nextCondLst stubs required by PowerPoint
    prev = etree.SubElement(seq, qn("p:prevCondLst"))
    prev_cond = etree.SubElement(prev, qn("p:cond"))
    prev_cond.set("evt", "onPrev")
    prev_cond.set("delay", "0")
    prev_tgt = etree.SubElement(prev_cond, qn("p:tgtEl"))
    etree.SubElement(prev_tgt, qn("p:sldTgt"))
    nxt = etree.SubElement(seq, qn("p:nextCondLst"))
    next_cond = etree.SubElement(nxt, qn("p:cond"))
    next_cond.set("evt", "onNext")
    next_cond.set("delay", "0")
    next_tgt = etree.SubElement(next_cond, qn("p:tgtEl"))
    etree.SubElement(next_tgt, qn("p:sldTgt"))


def hide_until_animated(shape):
    """Пометить фигуру скрытой до анимации (через nvPr)."""
    nvSpPr = shape._element.find(qn("p:nvSpPr"))
    if nvSpPr is None:
        return
    nvPr = nvSpPr.find(qn("p:nvPr"))
    if nvPr is None:
        nvPr = etree.SubElement(nvSpPr, qn("p:nvPr"))
    # p:ph not needed; use extLst timing? For appear, PowerPoint uses cNvPr invisible via set anim.
    # Also set a:effectLst empty; practical approach: start with transparency via solid fill alpha? Skip.
    pass


def etymon_slide(prs, n, ru, fr, boxes):
    slide = paper_slide(prs)
    eyebrow(slide, f"{n} / 5")
    word = add_text(slide, Inches(0.7), Inches(1.6), Inches(5.8), Inches(1.3), ru, 54, True, False, BURG, "Georgia")
    origin = add_text(slide, Inches(0.7), Inches(3.0), Inches(5.8), Inches(0.6), fr, 22, False, True, FR, "Georgia")
    anim = [word, origin]
    y = Inches(1.3)
    for title, body in boxes:
        c = card(slide, Inches(6.9), y, Inches(5.7), Inches(1.55))
        t1 = add_text(slide, Inches(7.1), y + Inches(0.12), Inches(5.3), Inches(0.3), title.upper(), 11, True, False, GOLD)
        t2 = add_text(slide, Inches(7.1), y + Inches(0.42), Inches(5.3), Inches(1.0), body, 16, False, False, INK)
        anim.append([c, t1, t2])
        y += Inches(1.75)
    add_click_appears(slide, anim)
    return slide


def build():
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H

    # 1 Title
    s = prs.slides.add_slide(prs.slide_layouts[6])
    rect(s, 0, 0, W / 2, H, RGBColor(0x3D, 0x5A, 0x80))
    rect(s, W / 2, 0, W / 2, H, RGBColor(0x8E, 0x24, 0x28))
    card(s, Inches(2.4), Inches(1.35), Inches(8.5), Inches(4.9))
    h1 = add_text(s, Inches(2.7), Inches(1.7), Inches(7.9), Inches(1.6),
                  "Лингвистический взлом: Русский язык через слова, которые вы уже знаете",
                  28, True, False, INK, "Georgia", PP_ALIGN.CENTER)
    fr = add_text(s, Inches(2.7), Inches(3.35), Inches(7.9), Inches(0.5),
                  "Le russe par les mots que vous connaissez déjà",
                  18, False, True, FR, "Georgia", PP_ALIGN.CENTER)
    sp = add_text(s, Inches(2.7), Inches(4.1), Inches(7.9), Inches(0.45),
                  "ШАГБАНОВА ХАБИБА САДЫРОВНА", 20, True, False, INK, "Georgia", PP_ALIGN.CENTER)
    role = add_text(s, Inches(2.9), Inches(4.65), Inches(7.5), Inches(1.2),
                    "доктор филологических наук, профессор, заместитель директора по социальной и культурной адаптации и межкультурной коммуникации ИРЯ РУДН им. Патриса Лумумбы",
                    13, False, False, MUTED, align=PP_ALIGN.CENTER)
    add_click_appears(s, [h1, fr, sp, role])

    # 2 Goal
    s = paper_slide(prs)
    eyebrow(s, "Слайд 2")
    heading(s, "Цель и задачи")
    c1 = card(s, Inches(0.7), Inches(1.6), Inches(5.3), Inches(5.1))
    t_goal = add_text(s, Inches(0.95), Inches(1.85), Inches(4.85), Inches(0.35), "ЦЕЛЬ", 12, True, False, BURG)
    b_goal = add_text(s, Inches(0.95), Inches(2.3), Inches(4.85), Inches(3.8),
                      "Демонстрация метода снятия языкового барьера у франкоговорящих студентов уровня A0–A1 через активацию пассивного лексического запаса.",
                      20, False, False, INK)
    tasks = [
        ("Детектив, не зубрёжка", "Трансформировать страх перед кириллицей в азарт исследователя."),
        ("Фонетические ловушки", "Отработать ударение, звук [ы] и букву Ё в знакомых корнях."),
        ("Ситуация успеха", "Студент понимает слова с первой минуты."),
    ]
    anim = [[c1, t_goal, b_goal]]
    y = Inches(1.6)
    for title, body in tasks:
        c = card(s, Inches(6.3), y, Inches(6.3), Inches(1.55))
        t = add_text(s, Inches(6.55), y + Inches(0.18), Inches(5.9), Inches(0.4), title, 20, True, False, INK, "Georgia")
        b = add_text(s, Inches(6.55), y + Inches(0.62), Inches(5.9), Inches(0.75), body, 15, False, False, MUTED)
        anim.append([c, t, b])
        y += Inches(1.7)
    add_click_appears(s, anim)

    # 3 Paris + Moscow on one slide (half + half, then optional emphasis)
    s = dark_slide(prs)
    eiffel = ROOT / "assets" / "morph-eiffel.jpg"
    basil = ROOT / "assets" / "morph-basil.jpg"
    half = W / 2
    if eiffel.exists():
        pic_l = s.shapes.add_picture(str(eiffel), 0, 0, half, H)
    else:
        pic_l = rect(s, 0, 0, half, H, FR)
        add_text(s, Inches(0.4), Inches(3.2), half - Inches(0.8), Inches(0.6),
                 "Париж", 32, True, False, WHITE, "Georgia", PP_ALIGN.CENTER)
    if basil.exists():
        pic_r = s.shapes.add_picture(str(basil), half, 0, half, H)
    else:
        pic_r = rect(s, half, 0, half, H, RU)
        add_text(s, half + Inches(0.4), Inches(3.2), half - Inches(0.8), Inches(0.6),
                 "Москва", 32, True, False, WHITE, "Georgia", PP_ALIGN.CENTER)
    # soft divider only — no city captions
    rect(s, half - Inches(0.03), 0, Inches(0.06), H, GOLD)

    # 4 Panic
    s = dark_slide(prs, RGBColor(0x1A, 0x08, 0x0A))
    k = add_text(s, Inches(0.8), Inches(1.0), Inches(11.7), Inches(0.5),
                 "CATASTROPHE · PANIQUE · QUE FAIRE ?", 18, True, False, GOLD, align=PP_ALIGN.CENTER)
    win = card(s, Inches(2.4), Inches(1.9), Inches(8.5), Inches(4.2), RGBColor(0x2A, 0x12, 0x14))
    bar = add_text(s, Inches(2.7), Inches(2.15), Inches(7.9), Inches(0.4),
                   "Language pack installer · v.1892", 14, True, False, CREAM)
    stars = add_text(s, Inches(2.7), Inches(3.0), Inches(7.9), Inches(0.6),
                     "* * * * * * * *", 28, True, False, WHITE, "Courier New", PP_ALIGN.CENTER)
    err1 = add_text(s, Inches(2.7), Inches(3.9), Inches(7.9), Inches(0.55),
                    "ERROR: LANGUAGE PACK NOT FOUND", 24, True, False, RGBColor(0xE0, 0x70, 0x70), "Courier New", PP_ALIGN.CENTER)
    err2 = add_text(s, Inches(2.7), Inches(4.7), Inches(7.9), Inches(0.55),
                    "SYSTEM FAILURE · ЧТО ДЕЛАТЬ?", 18, True, False, RGBColor(0xE0, 0x70, 0x70), align=PP_ALIGN.CENTER)
    add_click_appears(s, [k, [win, bar], stars, err1, err2])

    # 5 Decode
    s = dark_slide(prs, RGBColor(0x2A, 0x0D, 0x10))
    words = []
    y = Inches(1.8)
    for w in ("АВАНТЮРА", "БУЛЬВАР", "ВОЯЖ"):
        stars = add_text(s, Inches(1), y, Inches(11.3), Inches(0.9),
                         "* * * * * * * *", 42, True, False, RGBColor(0xE0, 0x70, 0x70), "Courier New", PP_ALIGN.CENTER)
        real = add_text(s, Inches(1), y, Inches(11.3), Inches(0.9),
                        w, 48, True, False, CREAM, "Courier New", PP_ALIGN.CENTER)
        words.append([stars, real])
        y += Inches(1.4)
    # click: show decoded words one by one (replace stars visually by appearing real on top)
    add_click_appears(s, [[w[1]] for w in words])

    # 6 Alphabet gallery (editable sample pairs)
    s = paper_slide(prs)
    eyebrow(s, "Алфавит через галлицизмы")
    heading(s, "Слова, которые вы уже знаете")
    pairs = [
        ("Авантюра", "aventure"), ("Бульвар", "boulevard"), ("Вояж", "voyage"),
        ("Гурман", "gourmand"), ("Дебют", "début"), ("Жалюзи", "jalousie"),
        ("Кошмар", "cauchemar"), ("Люстра", "lustre"), ("Пальто", "paletot"),
        ("Шедевр", "chef-d'œuvre"), ("Этикет", "étiquette"), ("Ювелир", "joaillier"),
    ]
    anim = []
    x, y = Inches(0.7), Inches(1.55)
    for i, (ru, frw) in enumerate(pairs):
        if i and i % 4 == 0:
            x = Inches(0.7)
            y += Inches(1.75)
        c = card(s, x, y, Inches(2.95), Inches(1.5))
        t1 = add_text(s, x + Inches(0.15), y + Inches(0.3), Inches(2.65), Inches(0.5), ru, 22, True, False, INK, "Georgia", PP_ALIGN.CENTER)
        t2 = add_text(s, x + Inches(0.15), y + Inches(0.85), Inches(2.65), Inches(0.4), frw, 16, False, True, FR, "Georgia", PP_ALIGN.CENTER)
        anim.append([c, t1, t2])
        x += Inches(3.15)
    note = add_text(s, Inches(0.7), Inches(6.85), Inches(12), Inches(0.35),
                    "В веб-версии слова идут бегущей строкой; здесь — ключевые примеры для разбора.",
                    12, False, False, MUTED)
    anim.append(note)
    add_click_appears(s, anim)

    # 7 Passport
    s = paper_slide(prs)
    eyebrow(s, "Certificat de nationalité russe")
    heading(s, "Русский паспорт французских слов")
    add_text(s, Inches(0.7), Inches(1.4), Inches(12), Inches(0.4),
             "Они живут по нашим законам и платят налоги ударениями.", 16, False, False, MUTED)
    anim = []
    entries = [
        ("Очередь · 1", "Бульон", "bouillon", "bouillon — бульон"),
        ("Очередь · 2", "Котлета", "côtelette", "côtelette — котлета"),
        ("Очередь · 3", "Винегрет", "vinaigrette", "vinaigrette — винегрет"),
    ]
    x = Inches(0.7)
    for q, ru, frw, line in entries:
        c = card(s, x, Inches(2.1), Inches(3.9), Inches(4.5))
        t0 = add_text(s, x + Inches(0.25), Inches(2.35), Inches(3.4), Inches(0.35), q.upper(), 12, True, False, BURG)
        t1 = add_text(s, x + Inches(0.25), Inches(3.0), Inches(3.4), Inches(0.7), ru, 36, True, False, BURG, "Georgia")
        t2 = add_text(s, x + Inches(0.25), Inches(3.8), Inches(3.4), Inches(0.4), frw, 18, False, True, FR, "Georgia")
        t3 = add_text(s, x + Inches(0.25), Inches(4.35), Inches(3.4), Inches(0.4), line, 14, False, True, FR)
        stamp = add_text(s, x + Inches(1.5), Inches(5.3), Inches(2.1), Inches(0.55), "ПРИНЯТО", 20, True, False, RGBColor(0x1B, 0x3F, 0xA0), align=PP_ALIGN.CENTER)
        stamp_box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x + Inches(1.55), Inches(5.25), Inches(2.0), Inches(0.6))
        stamp_box.fill.background()
        stamp_box.line.color.rgb = RGBColor(0x1B, 0x3F, 0xA0)
        stamp_box.line.width = Pt(2.5)
        anim.append([c, t0, t1, t2, t3])
        anim.append([stamp, stamp_box])
        x += Inches(4.15)
    add_click_appears(s, anim)

    # 8 False friends intro
    s = paper_slide(prs)
    eyebrow(s, "Ложные друзья переводчика")
    heading(s, "Не копия, а превращение")
    lead = add_text(s, Inches(0.7), Inches(1.55), Inches(11.8), Inches(1.0),
                    "Пять слов, в которых французский корень остался, а русская «маска» изменила форму, звук или оттенок смысла.",
                    20, False, False, MUTED)
    anim = [lead]
    x = Inches(1.2)
    for w in ("Шарлатан", "Шезлонг", "Кошмар", "Кашне", "Кашпо"):
        ch, width = chip(s, w, x, Inches(3.2))
        anim.append(ch)
        x += width + Inches(0.2)
    add_click_appears(s, anim)

    # 9–13 etymons
    etymon_slide(prs, 1, "Шарлатан", "Charlatan  [шарлатА]", [
        ("Во Франции", "Уличный торговец зельями, болтун. От ит. ciarlare — «болтать»."),
        ("В русском", "Обманщик, мошенник, выдающий себя за знатока."),
        ("Тип", "Прямое заимствование. Написание почти без изменений."),
    ])
    etymon_slide(prs, 2, "Шезлонг", "Chaise longue  [шэз лонг]", [
        ("Дословно", "Стул (chaise) длинный (longue)."),
        ("Трансформация", "Два слова слились в одно существительное мужского рода."),
        ("Тип", "Слияние. Значение мебели для полулежачего отдыха сохранилось."),
    ])
    etymon_slide(prs, 3, "Кошмар", "Cauchemar  [кошмАр]", [
        ("Состав", "Coucher — укладывать, давить. Mar — злой дух, демон."),
        ("Дословно", "«Демон, который давит на грудь во сне»."),
        ("В русском", "И страшный сон, и любое тяжёлое событие: «политический кошмар»."),
    ])
    etymon_slide(prs, 4, "Кашне", "Cache-nez  [кашнЭ]", [
        ("Дословно", "Спрячь нос! Cacher — прятать + nez — нос."),
        ("Трансформация", "Исчезли дефис и французский апостроф. Слово стало цельным."),
        ("Смысл", "Шейный платок. Сначала — стиль и защита лица, не только тепло."),
    ])
    etymon_slide(prs, 5, "Кашпо", "Cache-pot  [кашпО]", [
        ("Дословно", "Спрячь горшок! Cacher — прятать + pot — горшок."),
        ("Пара к кашне", "Та же модель: две части срослись по правилам русской грамматики."),
        ("Смысл", "Декоративная оболочка, скрывающая неприглядный цветочный горшок."),
    ])

    # 14 Table
    s = paper_slide(prs)
    eyebrow(s, "Слайд 4 · Алфавитный шок")
    heading(s, "Что мы уже можем")
    rows = [
        ("Русское слово", "Французский оригинал", "Дословный перевод", "Тип изменения"),
        ("Шарлатан", "Charlatan", "Болтун / мошенник", "Прямое заимствование"),
        ("Шезлонг", "Chaise longue", "Стул длинный", "Слияние двух слов"),
        ("Кошмар", "Cauchemar", "Укладывающий демон", "Фонетическое заимствование"),
        ("Кашне", "Cache-nez", "Спрячь нос", "Калька (сложное слово)"),
        ("Кашпо", "Cache-pot", "Спрячь горшок", "Калька (сложное слово)"),
    ]
    anim = []
    y = Inches(1.55)
    widths = [Inches(2.6), Inches(3.0), Inches(3.4), Inches(3.4)]
    for i, row in enumerate(rows):
        x = Inches(0.7)
        group = []
        for j, cell in enumerate(row):
            bg = GOLD if i == 0 else (WHITE if i % 2 else RGBColor(0xFA, 0xF4, 0xEA))
            ink = INK if i else NAVY
            r = rect(s, x, y, widths[j] - Inches(0.05), Inches(0.72), bg)
            t = add_text(s, x + Inches(0.1), y + Inches(0.15), widths[j] - Inches(0.25), Inches(0.45),
                         cell, 14 if i else 13, i == 0, False, ink)
            group.extend([r, t])
            x += widths[j]
        anim.append(group)
        y += Inches(0.78)
    add_click_appears(s, anim)

    # 15 Suitcase game
    s = paper_slide(prs)
    eyebrow(s, "Слайд 5 · Практикум")
    heading(s, "Кто быстрее соберёт свой саквояж")
    express = add_text(s, Inches(0.7), Inches(1.4), Inches(12), Inches(0.45),
                       "Экспресс «Африка — Россия»", 26, True, False, BURG, "Georgia", PP_ALIGN.CENTER)
    left = ["Балет", "Батон", "Берет", "Браслет", "Бюро", "Гардероб", "Гурман", "Жалюзи", "Жакет", "Кашне"]
    right = ["Колье", "Корсет", "Крем", "Люстра", "Маникюр", "Пальто", "Парфюм", "Помада", "Портфель", "Туалет", "Фуршет"]
    left_shapes, right_shapes = [], []
    bag_items(s, left, Inches(0.5), Inches(2.1), max_r=Inches(4.5), shapes_out=left_shapes)
    suitcase = card(s, Inches(5.0), Inches(2.6), Inches(3.3), Inches(3.2), RGBColor(0x5A, 0x2A, 0x18))
    suitcase.line.color.rgb = GOLD
    add_text(s, Inches(5.2), Inches(3.6), Inches(2.9), Inches(0.6), "САКВОЯЖ", 22, True, False, CREAM, "Georgia", PP_ALIGN.CENTER)
    lamp = add_text(s, Inches(5.3), Inches(5.1), Inches(2.7), Inches(0.45), "ПРОХОДИТЕ", 14, True, False, RGBColor(0xE8, 0xFF, 0xE8), align=PP_ALIGN.CENTER)
    lamp_bg = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.55), Inches(5.05), Inches(2.2), Inches(0.5))
    fill(lamp_bg, RGBColor(0x1F, 0x6B, 0x45))
    bag_items(s, right, Inches(8.6), Inches(2.1), max_r=Inches(12.8), shapes_out=right_shapes)
    add_click_appears(s, [express, left_shapes, [suitcase, lamp, lamp_bg], right_shapes])

    # 16 Reading + reveal extras
    s = paper_slide(prs)
    eyebrow(s, "Слайд 6 · Чтение")
    heading(s, "Кто лучше прочитает")
    c_ok = card(s, Inches(0.5), Inches(1.5), Inches(6.0), Inches(5.4))
    h_ok = add_text(s, Inches(0.75), Inches(1.7), Inches(5.5), Inches(0.45), "В саквояже", 24, True, False, INK, "Georgia")
    ok = ["Батон", "Берет", "Браслет", "Гардероб", "Жакет", "Кашне", "Колье", "Корсет", "Крем", "Пальто", "Парфюм", "Помада", "Портфель"]
    ok_shapes = []
    bag_items(s, ok, Inches(0.75), Inches(2.35), max_r=Inches(6.2), shapes_out=ok_shapes)
    c_ex = card(s, Inches(6.8), Inches(1.5), Inches(5.9), Inches(5.4))
    h_ex = add_text(s, Inches(7.05), Inches(1.7), Inches(5.4), Inches(0.45), "Какие слова лишние?", 22, True, False, INK, "Georgia")
    extras = ["Фуршет", "Туалет", "Балет", "Бюро", "Гурман", "Жалюзи", "Люстра", "Маникюр"]
    ex_shapes = []
    bag_items(s, extras, Inches(7.05), Inches(2.5), extra=True, max_r=Inches(12.5), shapes_out=ex_shapes)
    add_click_appears(s, [[c_ok, h_ok] + ok_shapes, [c_ex, h_ex], ex_shapes])

    # 17 Songs / prize
    s = paper_slide(prs)
    eyebrow(s, "Слайд 7 · Кто самый внимательный")
    heading(s, "Услышьте знакомые слова")
    prize = card(s, Inches(0.7), Inches(1.5), Inches(11.9), Inches(1.3))
    prize.line.color.rgb = GOLD
    ptxt = add_text(s, Inches(1.0), Inches(1.75), Inches(11.3), Inches(0.85),
                    "А сейчас тому, кто услышит знакомые слова, — приз: моё учебное пособие «Русский Экспресс».",
                    18, False, False, BURG)
    c1 = card(s, Inches(0.7), Inches(3.2), Inches(5.8), Inches(3.4))
    t1 = add_text(s, Inches(1.0), Inches(3.45), Inches(5.3), Inches(0.35), "CHEVALIER DE LA TABLE RONDE", 12, True, False, BURG)
    b1 = add_text(s, Inches(1.0), Inches(4.0), Inches(5.3), Inches(2.2),
                  "Chevalier de la table ronde\nDites-moi s’il le vin est bon",
                  20, False, False, INK, "Georgia")
    c2 = card(s, Inches(6.8), Inches(3.2), Inches(5.8), Inches(3.4))
    t2 = add_text(s, Inches(7.1), Inches(3.45), Inches(5.3), Inches(0.35), "MON MEC À MOI", 12, True, False, BURG)
    b2 = add_text(s, Inches(7.1), Inches(4.0), Inches(5.3), Inches(2.2),
                  "Mon mec à moi il me parle d’aventures\nEt quand il brille dans ses yeux\nj’pourrai passer la nuit",
                  18, False, False, INK, "Georgia")
    add_click_appears(s, [[prize, ptxt], [c1, t1, b1], [c2, t2, b2]])

    # 18 Hugo
    s = paper_slide(prs)
    eyebrow(s, "Виктор Гюго · «Париж»")
    heading(s, "Корни, которые уже ваши")
    stanzas = [
        "Cette ville aux longs cris,\nQui profile son front gris,\nDes toits frêles, cent tourelles,\nClochers grêles,\nC’est Paris.",
        "Des quadrilles,\nDes chansons\nMêlent filles et garçons.\nQuelle fête\nQue des têtes\nSur les faîtes\nDes maisons !",
        "Le vieux Louvre,\nLarge et lourd,\nIl ne s’ouvre\nQu’au grand jour.\n…\nComme une onde\nSur la mer.",
    ]
    anim = []
    x = Inches(0.6)
    for st in stanzas:
        c = card(s, x, Inches(1.55), Inches(4.0), Inches(5.3))
        t = add_text(s, x + Inches(0.25), Inches(1.9), Inches(3.5), Inches(4.6), st, 16, False, False, INK, "Georgia")
        anim.append([c, t])
        x += Inches(4.2)
    add_click_appears(s, anim)

    # 19 Map (editable labels)
    s = dark_slide(prs)
    eyebrow(s, "Геополитика смыслов", color=GOLD)
    africa = rect(s, Inches(0.6), Inches(1.3), Inches(5.4), Inches(5.5), RGBColor(0x3D, 0x5A, 0x80))
    russia = rect(s, Inches(7.0), Inches(1.5), Inches(5.6), Inches(5.1), RGBColor(0x8E, 0x24, 0x28))
    a_lab = add_text(s, Inches(0.8), Inches(1.5), Inches(5.0), Inches(0.4), "AFRIQUE FRANCOPHONE", 14, True, False, WHITE)
    r_lab = add_text(s, Inches(7.3), Inches(3.5), Inches(5.0), Inches(0.7), "РОССИЯ", 36, True, False, WHITE, "Georgia", PP_ALIGN.CENTER)
    cities = "Алжир · Рабат · Дакар · Бамако · Ниамей · Абиджан · Ломе · Яунде · Банги · Киншаса · Каир · Антананариву"
    c_txt = add_text(s, Inches(0.9), Inches(5.6), Inches(5.0), Inches(0.9), cities, 13, False, False, WHITE)
    rays = add_text(s, Inches(5.5), Inches(3.6), Inches(2.2), Inches(0.5), "······→", 22, True, False, GOLD, align=PP_ALIGN.CENTER)
    note = add_text(s, Inches(7.3), Inches(6.2), Inches(5.0), Inches(0.4),
                    "Нити смыслов: из столиц франкофонной Африки — в Россию", 12, False, False, CREAM, align=PP_ALIGN.CENTER)
    add_click_appears(s, [[africa, a_lab, c_txt], [russia, r_lab], [rays, note]])

    # 20 Finale
    s = dark_slide(prs, RGBColor(0x12, 0x0E, 0x0B))
    left = rect(s, 0, 0, W / 2, H, RGBColor(0x3D, 0x5A, 0x80))
    right = rect(s, W / 2, 0, W / 2, H, RGBColor(0x8E, 0x24, 0x28))
    num = add_text(s, Inches(0.8), Inches(1.3), Inches(11.7), Inches(1.0), "30", 72, True, False, GOLD, "Georgia", PP_ALIGN.CENTER)
    title = add_text(s, Inches(1.2), Inches(2.5), Inches(10.9), Inches(1.5),
                     "Русский язык ближе,\nчем вам кажется", 40, True, False, WHITE, "Georgia", PP_ALIGN.CENTER)
    cities = add_text(s, Inches(1.5), Inches(4.4), Inches(10.3), Inches(1.1),
                      "Бамако · Ниамей · Банги\nАбиджан · Яунде · Того\nи Россия",
                      20, False, False, CREAM, align=PP_ALIGN.CENTER)
    merci = add_text(s, Inches(1.5), Inches(5.9), Inches(10.3), Inches(0.6),
                     "Merci pour votre attention !", 24, False, True, WHITE, "Georgia", PP_ALIGN.CENTER)
    add_click_appears(s, [num, title, cities, merci])

    prs.save(OUT)
    print(f"OK → {OUT}")
    print(f"Слайдов: {len(prs.slides)}")


if __name__ == "__main__":
    build()
